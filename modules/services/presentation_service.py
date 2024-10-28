# services/presentation_service.py

import logging
import threading
from typing import Any, Dict, List, Optional
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from modules.utilities.logging_manager import setup_logging
from modules.utilities.config_loader import ConfigLoader
from modules.security.encryption_manager import EncryptionManager
from modules.security.authentication import AuthenticationManager


class PresentationServiceError(Exception):
    """Custom exception for PresentationService-related errors."""
    pass


class PresentationService:
    """
    Provides presentation creation and manipulation capabilities, including creating slides,
    adding text, images, charts, and formatting. Utilizes the python-pptx library to ensure
    comprehensive presentation handling. Ensures secure handling of files and configurations.
    """

    def __init__(self):
        """
        Initializes the PresentationService with necessary configurations and authentication.
        """
        self.logger = setup_logging('PresentationService')
        self.config_loader = ConfigLoader()
        self.encryption_manager = EncryptionManager()
        self.auth_manager = AuthenticationManager()
        self.lock = threading.Lock()
        self.logger.info("PresentationService initialized successfully.")

    def create_presentation(self, output_path: str, title_slide: Dict[str, str], slides: List[Dict[str, Any]]) -> bool:
        """
        Creates a new PowerPoint presentation with specified slides.

        Args:
            output_path (str): The file path for the created presentation.
            title_slide (Dict[str, str]): A dictionary containing 'title' and 'subtitle' for the title slide.
            slides (List[Dict[str, Any]]): A list of dictionaries, each representing a slide with its content.

        Returns:
            bool: True if the presentation is created successfully, False otherwise.
        """
        try:
            self.logger.debug(f"Creating presentation at '{output_path}' with {len(slides)} slides.")
            with self.lock:
                prs = Presentation()

                # Add Title Slide
                title_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(title_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                title.text = title_slide.get('title', 'Title')
                subtitle.text = title_slide.get('subtitle', 'Subtitle')
                self.logger.debug("Title slide added.")

                # Add Content Slides
                for idx, slide_content in enumerate(slides, start=1):
                    layout = prs.slide_layouts[1]  # Title and Content
                    slide = prs.slides.add_slide(layout)
                    slide_title = slide.shapes.title
                    slide_title.text = slide_content.get('title', f'Slide {idx}')

                    content = slide.placeholders[1].text = slide_content.get('content', '')

                    # Add Images
                    images = slide_content.get('images', [])
                    for image in images:
                        if os.path.exists(image.get('path', '')):
                            left = Inches(image.get('left', 1))
                            top = Inches(image.get('top', 2))
                            height = Inches(image.get('height', 2))
                            slide.shapes.add_picture(image['path'], left, top, height=height)
                            self.logger.debug(f"Image '{image['path']}' added to slide '{slide_title.text}'.")
                        else:
                            self.logger.warning(f"Image file '{image.get('path', '')}' does not exist. Skipping.")

                    # Add Charts
                    charts = slide_content.get('charts', [])
                    for chart in charts:
                        chart_type = chart.get('type', 'column')
                        data = chart.get('data', {})
                        categories = data.get('categories', [])
                        values = data.get('values', [])
                        chart_title = chart.get('title', 'Chart Title')

                        x, y, cx, cy = Inches(5), Inches(2), Inches(4), Inches(3)
                        chart_shape = slide.shapes.add_chart(
                            XL_CHART_TYPE[chart_type.upper()],
                            x, y, cx, cy,
                            chart_data=self._create_chart_data(categories, values)
                        )
                        chart_shape.chart.has_title = True
                        chart_shape.chart.chart_title.text = chart_title
                        self.logger.debug(f"Chart '{chart_title}' added to slide '{slide_title.text}'.")

                prs.save(output_path)
            self.logger.info(f"Presentation created successfully at '{output_path}'.")
            return True
        except Exception as e:
            self.logger.error(f"Error creating presentation at '{output_path}': {e}", exc_info=True)
            return False

    def _create_chart_data(self, categories: List[str], values: List[float]) -> CategoryChartData:
        """
        Creates chart data for inclusion in a presentation.

        Args:
            categories (List[str]): The categories for the chart.
            values (List[float]): The corresponding values for each category.

        Returns:
            CategoryChartData: The chart data object.
        """
        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series('Series 1', values)
        self.logger.debug("Chart data created.")
        return chart_data

    def add_textbox(self, presentation_path: str, slide_index: int, text: str, left: float, top: float, width: float, height: float, font_size: int = 18, bold: bool = False, italic: bool = False, color: Optional[RGBColor] = None) -> bool:
        """
        Adds a textbox to a specified slide in an existing presentation.

        Args:
            presentation_path (str): The file path to the presentation.
            slide_index (int): The index of the slide to add the textbox (0-based).
            text (str): The text content of the textbox.
            left (float): The left position in inches.
            top (float): The top position in inches.
            width (float): The width of the textbox in inches.
            height (float): The height of the textbox in inches.
            font_size (int, optional): The font size of the text. Defaults to 18.
            bold (bool, optional): Whether the text is bold. Defaults to False.
            italic (bool, optional): Whether the text is italic. Defaults to False.
            color (Optional[RGBColor], optional): The color of the text. Defaults to None.

        Returns:
            bool: True if the textbox is added successfully, False otherwise.
        """
        try:
            self.logger.debug(f"Adding textbox to presentation '{presentation_path}' on slide {slide_index} with text '{text}'.")
            with self.lock:
                prs = Presentation(presentation_path)
                if slide_index < 0 or slide_index >= len(prs.slides):
                    self.logger.error(f"Slide index {slide_index} is out of range for presentation '{presentation_path}'.")
                    return False
                slide = prs.slides[slide_index]
                left_in = Inches(left)
                top_in = Inches(top)
                width_in = Inches(width)
                height_in = Inches(height)
                textbox = slide.shapes.add_textbox(left_in, top_in, width_in, height_in)
                text_frame = textbox.text_frame
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = text
                run.font.size = Pt(font_size)
                run.font.bold = bold
                run.font.italic = italic
                if color:
                    run.font.color.rgb = color
                self.logger.debug("Textbox added successfully.")
                prs.save(presentation_path)
            self.logger.info(f"Textbox added successfully to slide {slide_index} in presentation '{presentation_path}'.")
            return True
        except Exception as e:
            self.logger.error(f"Error adding textbox to presentation '{presentation_path}' on slide {slide_index}: {e}", exc_info=True)
            return False

    def add_shape(self, presentation_path: str, slide_index: int, shape_type: str, left: float, top: float, width: float, height: float, fill_color: Optional[RGBColor] = None, line_color: Optional[RGBColor] = None) -> bool:
        """
        Adds a shape to a specified slide in an existing presentation.

        Args:
            presentation_path (str): The file path to the presentation.
            slide_index (int): The index of the slide to add the shape (0-based).
            shape_type (str): The type of shape to add (e.g., 'rectangle', 'circle').
            left (float): The left position in inches.
            top (float): The top position in inches.
            width (float): The width of the shape in inches.
            height (float): The height of the shape in inches.
            fill_color (Optional[RGBColor], optional): The fill color of the shape. Defaults to None.
            line_color (Optional[RGBColor], optional): The line color of the shape. Defaults to None.

        Returns:
            bool: True if the shape is added successfully, False otherwise.
        """
        try:
            self.logger.debug(f"Adding shape '{shape_type}' to presentation '{presentation_path}' on slide {slide_index}.")
            with self.lock:
                prs = Presentation(presentation_path)
                if slide_index < 0 or slide_index >= len(prs.slides):
                    self.logger.error(f"Slide index {slide_index} is out of range for presentation '{presentation_path}'.")
                    return False
                slide = prs.slides[slide_index]
                left_in = Inches(left)
                top_in = Inches(top)
                width_in = Inches(width)
                height_in = Inches(height)
                shape = None

                shape_mapping = {
                    'rectangle': MSO_SHAPE.RECTANGLE,
                    'circle': MSO_SHAPE.OVAL,
                    'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
                    'diamond': MSO_SHAPE.DIAMOND,
                    'star': MSO_SHAPE.STAR_5,
                    'arrow': MSO_SHAPE.ARROW
                }

                if shape_type.lower() in shape_mapping:
                    shape = slide.shapes.add_shape(shape_mapping[shape_type.lower()], left_in, top_in, width_in, height_in)
                else:
                    self.logger.error(f"Unsupported shape type '{shape_type}'. Supported types: {list(shape_mapping.keys())}.")
                    return False

                if fill_color:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = fill_color
                    self.logger.debug(f"Shape fill color set to '{fill_color}'.")
                if line_color:
                    shape.line.color.rgb = line_color
                    self.logger.debug(f"Shape line color set to '{line_color}'.")

                prs.save(presentation_path)
            self.logger.info(f"Shape '{shape_type}' added successfully to slide {slide_index} in presentation '{presentation_path}'.")
            return True
        except Exception as e:
            self.logger.error(f"Error adding shape to presentation '{presentation_path}' on slide {slide_index}: {e}", exc_info=True)
            return False

    def add_table(self, presentation_path: str, slide_index: int, rows: int, cols: int, left: float, top: float, width: float, height: float, data: Optional[List[List[Any]]] = None, style: Optional[str] = 'Light List Accent 1') -> bool:
        """
        Adds a table to a specified slide in an existing presentation.

        Args:
            presentation_path (str): The file path to the presentation.
            slide_index (int): The index of the slide to add the table (0-based).
            rows (int): Number of rows in the table.
            cols (int): Number of columns in the table.
            left (float): The left position in inches.
            top (float): The top position in inches.
            width (float): The width of the table in inches.
            height (float): The height of the table in inches.
            data (Optional[List[List[Any]]], optional): The data to populate the table. Defaults to None.
            style (Optional[str], optional): The style of the table. Defaults to 'Light List Accent 1'.

        Returns:
            bool: True if the table is added successfully, False otherwise.
        """
        try:
            self.logger.debug(f"Adding table to presentation '{presentation_path}' on slide {slide_index} with {rows} rows and {cols} columns.")
            with self.lock:
                prs = Presentation(presentation_path)
                if slide_index < 0 or slide_index >= len(prs.slides):
                    self.logger.error(f"Slide index {slide_index} is out of range for presentation '{presentation_path}'.")
                    return False
                slide = prs.slides[slide_index]
                left_in = Inches(left)
                top_in = Inches(top)
                width_in = Inches(width)
                height_in = Inches(height)
                table_shape = slide.shapes.add_table(rows, cols, left_in, top_in, width_in, height_in)
                table = table_shape.table

                if data:
                    for row_idx, row_data in enumerate(data):
                        for col_idx, cell_data in enumerate(row_data):
                            if row_idx < rows and col_idx < cols:
                                table.cell(row_idx, col_idx).text = str(cell_data)
                                self.logger.debug(f"Set cell ({row_idx}, {col_idx}) to '{cell_data}'.")
                table.style = style
                self.logger.debug(f"Table style set to '{style}'.")
                prs.save(presentation_path)
            self.logger.info(f"Table added successfully to slide {slide_index} in presentation '{presentation_path}'.")
            return True
        except Exception as e:
            self.logger.error(f"Error adding table to presentation '{presentation_path}' on slide {slide_index}: {e}", exc_info=True)
            return False

    def add_slide_layout(self, presentation_path: str, layout_title: str, layout_content: str, position: Optional[int] = None) -> bool:
        """
        Adds a custom slide layout to an existing presentation.

        Args:
            presentation_path (str): The file path to the presentation.
            layout_title (str): The title of the new slide layout.
            layout_content (str): The content placeholder text for the new slide layout.
            position (Optional[int], optional): The position to insert the new layout. Defaults to None.

        Returns:
            bool: True if the slide layout is added successfully, False otherwise.
        """
        try:
            self.logger.debug(f"Adding custom slide layout '{layout_title}' to presentation '{presentation_path}'.")
            # Note: python-pptx does not support adding new slide layouts programmatically.
            # This functionality would require manipulating the XML directly or using a template with predefined layouts.
            # Therefore, we'll log an error indicating the limitation.
            self.logger.error("Adding custom slide layouts is not supported by python-pptx library.")
            return False
        except Exception as e:
            self.logger.error(f"Error adding slide layout to presentation '{presentation_path}': {e}", exc_info=True)
            return False

    def close_service(self):
        """
        Closes any resources or sessions held by the service.
        """
        try:
            self.logger.debug("Closing PresentationService resources.")
            # Currently, no persistent resources to close
            self.logger.info("PresentationService closed successfully.")
        except Exception as e:
            self.logger.error(f"Error closing PresentationService: {e}", exc_info=True)
            raise PresentationServiceError(f"Error closing PresentationService: {e}")
